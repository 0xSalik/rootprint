"""Hunarmand 7-slide editorial deck.

Run from this folder:

    pip install -r requirements.txt
    python build_deck.py

Outputs ``hunarmand.pptx`` next to this script. The deck is rendered
through python-pptx primitives only: rectangles, lines, ellipses, and
text boxes. Every shape sits on raw parchment; no decorative chrome,
no AI residue.

Typography is JetBrains Mono everywhere (Light, Regular, Bold). The
TTFs are bundled in ``fonts/`` so the deck renders identically on a
reviewer's machine after a one-time install.
"""

from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree


# ── Palette (the website's, exact hex) ─────────────────────────────────
PARCHMENT       = RGBColor(0xFA, 0xF7, 0xF2)
AGED_PAPER      = RGBColor(0xF0, 0xEB, 0xE3)
WALNUT          = RGBColor(0x1C, 0x14, 0x10)
FADED_INK       = RGBColor(0x5A, 0x4A, 0x3A)
MUTED           = RGBColor(0x9A, 0x8A, 0x7A)
ROUGE           = RGBColor(0x8B, 0x1A, 0x1A)
ROUGE_HOVER     = RGBColor(0xC4, 0x45, 0x4A)
GOLD            = RGBColor(0xC8, 0x97, 0x5A)
EDGE            = RGBColor(0xD4, 0xC8, 0xB8)

# Pre-blended walnut over parchment for "alpha" effects (python-pptx
# doesn't expose text alpha through the public API in a portable way,
# so we precompute the visual colour we want).
def _blend(fg: RGBColor, bg: RGBColor, alpha: float) -> RGBColor:
    fr, fg_, fb = fg[0], fg[1], fg[2]
    br, bg_b, bb = bg[0], bg[1], bg[2]
    return RGBColor(
        int(round(alpha * fr + (1 - alpha) * br)),
        int(round(alpha * fg_ + (1 - alpha) * bg_b)),
        int(round(alpha * fb + (1 - alpha) * bb)),
    )

WALNUT_60 = _blend(WALNUT, PARCHMENT, 0.60)   # ~ #746e6a
WALNUT_50 = _blend(WALNUT, PARCHMENT, 0.50)   # ~ #8b8581
WALNUT_40 = _blend(WALNUT, PARCHMENT, 0.40)   # ~ #a19c97
WALNUT_30 = _blend(WALNUT, PARCHMENT, 0.30)   # ~ #b8b3ad
WALNUT_20 = _blend(WALNUT, PARCHMENT, 0.20)   # ~ #cfcbc4
FADED_35  = _blend(FADED_INK, PARCHMENT, 0.35)
FADED_60  = _blend(FADED_INK, PARCHMENT, 0.60)


# ── Type faces ─────────────────────────────────────────────────────────
JBM_REG   = "JetBrains Mono"
JBM_BOLD  = "JetBrains Mono"     # bold = font.bold = True
JBM_LIGHT = "JetBrains Mono Light"  # the Light face is its own family
JBM_ITAL  = "JetBrains Mono"     # italic via font.italic


# ── XML helpers (python-pptx public API gaps) ──────────────────────────
def set_char_spacing(run, hundredths_of_pt: int) -> None:
    """Letter-spacing via OOXML rPr/spc (1/100 of a point)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(hundredths_of_pt))


def set_text_alpha(run, alpha_pct: int) -> None:
    """Set transparency on a text run's solid fill via OOXML.

    python-pptx doesn't expose alpha on solid fills directly, so we drop
    into the rPr/solidFill/srgbClr/alpha element. The fill must already
    be set (via run.font.color.rgb = ...) before this is called.
    """
    rPr = run._r.get_or_add_rPr()
    solid = rPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Remove any existing <a:alpha/>
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(alpha_pct * 1000)))  # OOXML: 0..100000


def disable_autosize(tf) -> None:
    """Stop python-pptx from shrinking text to fit the box."""
    bodyPr = tf._txBody.bodyPr
    for child in list(bodyPr):
        if child.tag in (qn("a:normAutofit"), qn("a:spAutoFit")):
            bodyPr.remove(child)
    # Ensure no auto-fit
    no_fit = etree.SubElement(bodyPr, qn("a:noAutofit"))


# ── Slide primitives ───────────────────────────────────────────────────
def add_parchment_bg(slide, width_in: float, height_in: float) -> None:
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(width_in), Inches(height_in)
    )
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = PARCHMENT
    rect.shadow.inherit = False


def add_text(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = JBM_REG,
    size: int = 11,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WALNUT,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    spacing_hpt: Optional[int] = None,
    line_spacing: Optional[float] = None,
) -> "Shape":
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    disable_autosize(tf)
    tf.vertical_anchor = anchor

    # Multi-line: split on newlines into paragraphs.
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if line_spacing is not None:
            para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if spacing_hpt is not None:
            set_char_spacing(run, spacing_hpt)
    return box


def add_line(
    slide,
    *,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: RGBColor = WALNUT,
    weight_pt: float = 0.75,
    dash: Optional[str] = None,  # "dash", "dot", "sysDash", "sysDot"
) -> "Shape":
    conn = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line = conn.line
    line.color.rgb = color
    line.width = Pt(weight_pt)
    if dash:
        ln = conn.line._get_or_add_ln()
        # Remove any existing dash elem
        for el in ln.findall(qn("a:prstDash")):
            ln.remove(el)
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", dash)
    return conn


def add_disk(
    slide,
    *,
    cx: float,
    cy: float,
    diameter_pt: float,
    fill: Optional[RGBColor] = WALNUT,
    line: Optional[RGBColor] = None,
    line_pt: float = 0.75,
) -> "Shape":
    d_in = diameter_pt / 72.0
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - d_in / 2),
        Inches(cy - d_in / 2),
        Inches(d_in),
        Inches(d_in),
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_pt)
    shp.shadow.inherit = False
    return shp


def add_filled_rect(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: Optional[RGBColor] = None,
    line_pt: float = 0.5,
) -> "Shape":
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_pt)
    shp.shadow.inherit = False
    return shp


# ── Footer (every slide) ───────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5


def add_footer(slide, n: int, total: int = 6) -> None:
    # Wordmark "hunarmand"
    add_text(
        slide,
        "hunarmand",
        left=SLIDE_W - 1.2,
        top=SLIDE_H - 0.55,
        width=1.0,
        height=0.18,
        font_name=JBM_LIGHT,
        size=8,
        color=WALNUT_40,
        align=PP_ALIGN.RIGHT,
        spacing_hpt=120,
    )
    # Slide number
    add_text(
        slide,
        f"{n:02d} / {total:02d}",
        left=SLIDE_W - 1.2,
        top=SLIDE_H - 0.32,
        width=1.0,
        height=0.18,
        font_name=JBM_LIGHT,
        size=8,
        color=WALNUT_40,
        align=PP_ALIGN.RIGHT,
        spacing_hpt=120,
    )


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 01 — cover
# ─────────────────────────────────────────────────────────────────────
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top-left orientation
    add_text(
        slide,
        "tacit knowledge os",
        left=0.6, top=0.55, width=4.0, height=0.3,
        font_name=JBM_LIGHT, size=14, color=FADED_INK,
    )
    add_text(
        slide,
        "for kashmiri heritage artisans",
        left=0.6, top=0.85, width=4.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=WALNUT_40,
        spacing_hpt=120,
    )

    # Wordmark, vertically centered, on left third. JetBrains Mono Bold
    # at 84pt fits "hunarmand" (9 chars) on a single line within the
    # left half of the canvas without wrapping.
    word_top = (SLIDE_H / 2) - 1.0
    add_text(
        slide,
        "hunarmand",
        left=0.55, top=word_top, width=8.6, height=1.4,
        font_name=JBM_REG, size=84, bold=True, color=WALNUT,
    )
    # Rouge hairline under wordmark
    line_y = word_top + 1.35
    add_line(
        slide,
        x1=0.6, y1=line_y, x2=0.6 + 1.4, y2=line_y,
        color=ROUGE, weight_pt=0.9,
    )

    # Tagline under the wordmark, set in faded ink with breathing room.
    add_text(
        slide,
        "we treat what the master never wrote down\nas captureable, queryable, ownable infrastructure.",
        left=0.6, top=line_y + 0.25, width=8.0, height=0.85,
        font_name=JBM_LIGHT, size=14, color=FADED_INK,
        line_spacing=1.45,
    )

    # Vertical talim motif on the right side — a small column of angled
    # walnut strokes that ties visually to the methodology slide and
    # gives the cover a craft-rooted ornament instead of corporate chrome.
    motif_x = SLIDE_W - 1.6
    motif_top = 1.55
    cell_h = 0.55
    cell_w = 0.34
    import random as _r
    _r.seed(11)
    for i in range(7):
        cy = motif_top + i * cell_h
        ax1 = motif_x - cell_w / 2 + _r.uniform(0.02, 0.06)
        ay1 = cy + _r.uniform(0.04, 0.10)
        ax2 = motif_x + _r.uniform(-0.02, 0.06)
        ay2 = cy + cell_h * 0.45 + _r.uniform(-0.04, 0.04)
        add_line(slide, x1=ax1, y1=ay1, x2=ax2, y2=ay2,
                 color=WALNUT, weight_pt=0.8)
        bx1 = motif_x + _r.uniform(-0.03, 0.05)
        by1 = cy + cell_h * 0.40 + _r.uniform(-0.03, 0.04)
        bx2 = motif_x + cell_w / 2 - _r.uniform(0.02, 0.07)
        by2 = cy + cell_h * 0.85 + _r.uniform(-0.03, 0.04)
        add_line(slide, x1=bx1, y1=by1, x2=bx2, y2=by2,
                 color=WALNUT, weight_pt=0.8)

    # Caption next to the motif, oriented as a margin gloss.
    add_text(
        slide,
        "talim, kanihama,\nc. 1820.  the script\nthe master kept.",
        left=motif_x - 2.4, top=motif_top + 0.6, width=2.2, height=1.0,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        align=PP_ALIGN.RIGHT, spacing_hpt=60, line_spacing=1.5,
    )

    # Team imprint, bottom-left. "bitwise" reads as the publisher mark
    # of the work, with srinagar / 2026 as the colophon line beneath.
    imprint_y = SLIDE_H - 1.1
    add_text(
        slide,
        "bitwise",
        left=0.6, top=imprint_y, width=3.0, height=0.4,
        font_name=JBM_REG, size=18, bold=True, color=WALNUT,
    )
    # Tiny gold rule under the team mark, mirroring the rouge under the
    # wordmark — the gold thread that runs through the website's palette.
    add_line(
        slide,
        x1=0.6, y1=imprint_y + 0.34,
        x2=0.6 + 0.7, y2=imprint_y + 0.34,
        color=GOLD, weight_pt=0.9,
    )
    add_text(
        slide,
        "srinagar  ·  2026  ·  v 0.1 hackathon build",
        left=0.6, top=imprint_y + 0.42, width=6.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    # Right metadata, anchored above the footer, right-aligned.
    meta_w = 5.6
    meta_left = SLIDE_W - 0.6 - meta_w
    add_text(
        slide,
        "neon pgvector  ·  hf spaces  ·  render  ·  vercel",
        left=meta_left, top=imprint_y + 0.42, width=meta_w, height=0.28,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        align=PP_ALIGN.RIGHT, spacing_hpt=120,
    )

    # Gold leaf, low-right, clear of the footer column.
    add_disk(
        slide,
        cx=SLIDE_W - 0.55, cy=SLIDE_H - 1.5,
        diameter_pt=8, fill=GOLD,
    )

    add_footer(slide, 1)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 02 — the disappearance
# ─────────────────────────────────────────────────────────────────────
def slide_disappearance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(
        slide,
        "loss / kashmiri heritage crafts / 2026",
        left=0.6, top=0.45, width=6.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    # Lineage tree (left half), anchored at x = 1.2"
    tree_x = 1.2
    nodes = [
        # (y, diameter_pt, fill, line_color, alpha (line on hollow), caption)
        (1.10, 10, WALNUT,    None,      "ustaad, b. 1953"),
        (1.95,  9, WALNUT,    None,      "son, taxi driver"),
        (2.80,  8, WALNUT,    None,      "nephew, learning"),
        (3.65,  8, ROUGE,     None,      "grandson, exam prep"),
        # hollow rings + dotted line below this point
        (4.50,  7, None,      WALNUT,    "no apprentice"),
        (5.35,  6, None,      WALNUT_30, "no apprentice"),
        (6.20,  4, FADED_35,  None,      "silence"),
    ]

    for i, (cy, diameter, fill, line_color, _) in enumerate(nodes):
        add_disk(
            slide,
            cx=tree_x, cy=cy,
            diameter_pt=diameter,
            fill=fill,
            line=line_color,
            line_pt=0.9,
        )

    # Connectors between nodes
    # solid for first 3 segments, dotted for the rest
    for i in range(len(nodes) - 1):
        y1 = nodes[i][0]
        y2 = nodes[i + 1][0]
        # leave breathing space around the disk
        d1_in = nodes[i][1] / 72.0
        d2_in = nodes[i + 1][1] / 72.0
        seg_y1 = y1 + d1_in / 2 + 0.02
        seg_y2 = y2 - d2_in / 2 - 0.02
        is_dotted = i >= 3
        add_line(
            slide,
            x1=tree_x, y1=seg_y1, x2=tree_x, y2=seg_y2,
            color=WALNUT if not is_dotted else WALNUT_60,
            weight_pt=0.75,
            dash="sysDot" if is_dotted else None,
        )

    # Captions next to each node
    for cy, _, _, _, caption in nodes:
        add_text(
            slide,
            caption,
            left=tree_x + 0.22, top=cy - 0.12, width=3.6, height=0.25,
            font_name=JBM_LIGHT, size=9, color=FADED_INK,
            spacing_hpt=60,
        )

    # Right side — the number. The big "12" sits at 200pt; the "yrs"
    # unit floats next to its baseline at a much smaller size, the way
    # an editorial display number is set with its unit as a tag.
    big_left = 5.6
    big_top = 1.0
    big_w = SLIDE_W - big_left - 0.7

    # Right-aligned composite. Compute approximate widths so "12" is the
    # last big character and "yrs" sits to its right.
    yrs_w = 1.2
    add_text(
        slide,
        "yrs",
        left=SLIDE_W - 0.7 - yrs_w, top=big_top + 1.05,
        width=yrs_w, height=0.55,
        font_name=JBM_LIGHT, size=28, color=FADED_INK,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
    )
    num_w = big_w - yrs_w - 0.15
    add_text(
        slide,
        "12",
        left=big_left, top=big_top, width=num_w, height=2.4,
        font_name=JBM_LIGHT, size=140, color=WALNUT,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
    )

    # Subhead under number, with breathing room so the "2" descender
    # never touches the caption baseline.
    sub_top = big_top + 2.55
    add_text(
        slide,
        "median time a kashmiri craft survives without an apprentice",
        left=big_left, top=sub_top, width=big_w, height=0.4,
        font_name=JBM_REG, size=14, bold=True, color=WALNUT,
        align=PP_ALIGN.RIGHT,
    )

    # Body two lines
    body_top = sub_top + 0.5
    add_text(
        slide,
        "talking to ten thousand workshops in the valley\nis not a sentimental project. it is a deadline.",
        left=big_left, top=body_top, width=big_w, height=0.85,
        font_name=JBM_REG, size=11, color=WALNUT,
        align=PP_ALIGN.RIGHT, line_spacing=1.4,
    )

    # Rouge hairline below body, right-aligned
    rouge_y = body_top + 1.0
    add_line(
        slide,
        x1=SLIDE_W - 0.7 - 2.6, y1=rouge_y,
        x2=SLIDE_W - 0.7,       y2=rouge_y,
        color=ROUGE, weight_pt=1.2,
    )

    add_footer(slide, 2)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 03 — the solution (4-layer system + value capture)
# ─────────────────────────────────────────────────────────────────────
def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(
        slide,
        "solution / four interlocking layers, anchored on cryptographic provenance",
        left=0.6, top=0.45, width=12.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    # Diagonal pipeline of the four layers, slightly compressed so the
    # economics bands fit underneath.
    nodes = [
        # (x, y, diameter_pt, fill, label, function)
        (1.2,  1.4,  10, WALNUT, "vault",
         "captures what the master never wrote down"),
        (4.2,  2.55, 12, ROUGE,  "sanad",
         "signs the piece against the master's own keypair"),
        (7.2,  3.7,  10, WALNUT, "ustaad",
         "sells the master's time, not the middleman's"),
        (10.2, 4.85, 10, WALNUT, "bazaar",
         "sells the object with provenance attached"),
    ]
    import math
    for i in range(len(nodes) - 1):
        x1, y1, d1, *_ = nodes[i]
        x2, y2, d2, *_ = nodes[i + 1]
        d1_in, d2_in = d1 / 72.0, d2 / 72.0
        dx, dy = x2 - x1, y2 - y1
        L = math.hypot(dx, dy)
        ux, uy = dx / L, dy / L
        sx, sy = x1 + ux * d1_in / 2, y1 + uy * d1_in / 2
        ex, ey = x2 - ux * d2_in / 2, y2 - uy * d2_in / 2
        add_line(slide, x1=sx, y1=sy, x2=ex, y2=ey,
                 color=WALNUT, weight_pt=0.9)
    for x, y, d, fill, label, fn in nodes:
        add_disk(slide, cx=x, cy=y, diameter_pt=d, fill=fill)
        label_x = x + (d / 72.0) / 2 + 0.22
        avail_w = max(2.0, SLIDE_W - label_x - 0.6)
        add_text(slide, label,
                 left=label_x, top=y - 0.32, width=avail_w, height=0.32,
                 font_name=JBM_REG, size=15, bold=True, color=WALNUT)
        add_text(slide, fn,
                 left=label_x, top=y - 0.04, width=avail_w, height=0.32,
                 font_name=JBM_LIGHT, size=10, color=FADED_INK,
                 spacing_hpt=60)

    # Keystone gloss, anchored to sanad with a faint dotted leader.
    sanad_x, sanad_y = nodes[1][0], nodes[1][1]
    keystone_left = 9.8
    keystone_y = 0.85
    add_line(slide, x1=keystone_left, y1=keystone_y,
             x2=keystone_left + 1.4, y2=keystone_y,
             color=ROUGE, weight_pt=1.0)
    add_text(slide, "the keystone",
             left=keystone_left, top=keystone_y + 0.08, width=2.6, height=0.32,
             font_name=JBM_REG, size=11, bold=True, color=WALNUT)
    add_text(slide,
             "ed25519 + rfc 8785 jcs.\na counterfeit shawl can be made.\na sanad signed by mohammad yusuf cannot.",
             left=keystone_left, top=keystone_y + 0.36, width=3.3, height=1.0,
             font_name=JBM_LIGHT, size=9, color=FADED_INK,
             spacing_hpt=60, line_spacing=1.5)
    add_line(slide,
             x1=keystone_left + 0.1, y1=keystone_y + 1.2,
             x2=sanad_x + 0.18, y2=sanad_y - 0.12,
             color=WALNUT_30, weight_pt=0.5, dash="sysDot")

    # Economics bands, lower-third of the slide. Anchored low enough
    # to clear the bazaar disk above and high enough that the bottom
    # caption never grazes the footer at y=6.95.
    bar_x = 1.8
    bar_w = 9.0
    bar_h = 0.28
    top1 = 5.4

    # today bar — captions sit BELOW for consistency with hunarmand bar
    add_text(slide, "today",
             left=0.6, top=top1 + (bar_h - 0.22) / 2, width=1.1, height=0.32,
             font_name=JBM_REG, size=10, bold=True, color=WALNUT)
    seg1_w = bar_w * 0.78
    add_filled_rect(slide, left=bar_x, top=top1, width=seg1_w, height=bar_h,
                    fill=WALNUT)
    add_filled_rect(slide, left=bar_x + seg1_w, top=top1,
                    width=bar_w - seg1_w, height=bar_h, fill=FADED_35)
    add_text(slide, "middleman chain  ·  70 to 85 percent",
             left=bar_x + 0.05, top=top1 + bar_h + 0.04, width=seg1_w, height=0.22,
             font_name=JBM_LIGHT, size=9, color=FADED_INK, spacing_hpt=60)
    add_text(slide, "artisan  ·  15 to 30 percent",
             left=bar_x + seg1_w + 0.05, top=top1 + bar_h + 0.04,
             width=bar_w - seg1_w + 1.6, height=0.22,
             font_name=JBM_LIGHT, size=9, color=FADED_INK, spacing_hpt=60)

    # hunarmand bar
    top2 = top1 + bar_h + 0.45
    add_text(slide, "hunarmand",
             left=0.6, top=top2 + (bar_h - 0.22) / 2, width=1.1, height=0.32,
             font_name=JBM_REG, size=10, bold=True, color=ROUGE)
    seg1b_w = bar_w * 0.08
    add_filled_rect(slide, left=bar_x, top=top2, width=seg1b_w, height=bar_h,
                    fill=ROUGE)
    add_filled_rect(slide, left=bar_x + seg1b_w, top=top2,
                    width=bar_w - seg1b_w, height=bar_h, fill=WALNUT)
    add_text(slide, "platform  ·  6 to 9 percent",
             left=bar_x - 1.6, top=top2 + bar_h + 0.04,
             width=1.6 + seg1b_w - 0.05, height=0.22,
             font_name=JBM_LIGHT, size=9, color=FADED_INK,
             align=PP_ALIGN.RIGHT, spacing_hpt=60)
    add_text(slide, "artisan  ·  91 to 94 percent",
             left=bar_x + seg1b_w + 0.15, top=top2 + bar_h + 0.04,
             width=bar_w - seg1b_w, height=0.22,
             font_name=JBM_LIGHT, size=9, color=FADED_INK, spacing_hpt=60)

    add_footer(slide, 3)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 04 — methodology (talim insight + 5 stages of the AI pipeline)
# ─────────────────────────────────────────────────────────────────────
def slide_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(slide,
             "methodology / how tacit knowledge becomes infrastructure",
             left=0.6, top=0.45, width=12.0, height=0.25,
             font_name=JBM_LIGHT, size=9, color=FADED_INK,
             spacing_hpt=120)

    # ── Left half: the talim insight (compact) ──
    # Glyph stack, 6 cells
    glyph_x = 1.4
    glyph_top = 1.1
    cell_h = 0.5
    cell_w = 0.32
    import random as _r
    _r.seed(7)
    for i in range(6):
        cy = glyph_top + i * cell_h
        ax1 = glyph_x - cell_w / 2 + _r.uniform(0.02, 0.06)
        ay1 = cy + _r.uniform(0.04, 0.10)
        ax2 = glyph_x + _r.uniform(-0.02, 0.06)
        ay2 = cy + cell_h * 0.45 + _r.uniform(-0.04, 0.04)
        add_line(slide, x1=ax1, y1=ay1, x2=ax2, y2=ay2,
                 color=WALNUT, weight_pt=0.8)
        bx1 = glyph_x + _r.uniform(-0.03, 0.05)
        by1 = cy + cell_h * 0.40 + _r.uniform(-0.03, 0.04)
        bx2 = glyph_x + cell_w / 2 - _r.uniform(0.02, 0.07)
        by2 = cy + cell_h * 0.85 + _r.uniform(-0.03, 0.04)
        add_line(slide, x1=bx1, y1=by1, x2=bx2, y2=by2,
                 color=WALNUT, weight_pt=0.8)

    add_text(slide, "talim,  kanihama,  c. 1820",
             left=glyph_x - 1.05, top=glyph_top - 0.32, width=2.2, height=0.3,
             font_name=JBM_LIGHT, size=9, color=FADED_INK,
             spacing_hpt=60)

    # Rouge hairline arc connecting glyph to insight body
    add_line(slide, x1=glyph_x + 0.1, y1=glyph_top + 6 * cell_h - 0.1,
             x2=2.6, y2=4.85,
             color=ROUGE, weight_pt=1.0)

    # Insight body, beneath the glyph. Frames what the methodology
    # *produces* without echoing the closing slide's quote.
    insight_top = 4.95
    add_text(slide,
             "twelve hours of speech.\nfour hundred fields.\none signed artefact.",
             left=0.6, top=insight_top, width=4.6, height=1.0,
             font_name=JBM_REG, size=14, bold=True, color=WALNUT,
             line_spacing=1.4)
    add_text(slide,
             "the master's tacit knowledge,\nstructured under their own keypair,\nowned by the master in perpetuity.",
             left=0.6, top=insight_top + 1.05, width=5.0, height=1.0,
             font_name=JBM_LIGHT, size=11, color=FADED_INK, line_spacing=1.5)

    # ── Right half: the 5-stage methodology pipeline ──
    pipe_left = 6.4
    pipe_top = 1.05
    row_h = 0.95
    stages = [
        ("01", "session",
         "master in own workshop, interview led in koshur"),
        ("02", "transcribe",
         "bhashini  >  ai4bharat  >  groq  >  whisper, ladder"),
        ("03", "extract",
         "multi-pass llm, strict json schema, repair loop"),
        ("04", "sign",
         "ed25519 over rfc 8785 jcs, one keypair per master"),
        ("05", "index",
         "pgvector embeddings, public api in under a second"),
    ]
    for i, (num, label, caption) in enumerate(stages):
        ry = pipe_top + i * row_h
        # number, in Light, big, faded ink for hierarchy
        add_text(slide, num,
                 left=pipe_left, top=ry, width=0.9, height=0.7,
                 font_name=JBM_LIGHT, size=36, color=FADED_60,
                 anchor=MSO_ANCHOR.TOP)
        # walnut hairline under each row
        add_line(slide, x1=pipe_left + 1.0, y1=ry + 0.1,
                 x2=SLIDE_W - 0.6, y2=ry + 0.1,
                 color=EDGE, weight_pt=0.5)
        # label
        add_text(slide, label,
                 left=pipe_left + 1.05, top=ry + 0.18, width=2.0, height=0.32,
                 font_name=JBM_REG, size=14, bold=True, color=WALNUT)
        # caption
        add_text(slide, caption,
                 left=pipe_left + 1.05, top=ry + 0.46, width=SLIDE_W - pipe_left - 1.6,
                 height=0.32,
                 font_name=JBM_LIGHT, size=10, color=FADED_INK, spacing_hpt=60)

    add_footer(slide, 4)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 05 — technology (deployed stack + live numbers)
# ─────────────────────────────────────────────────────────────────────
def slide_technology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(slide,
             "technology / deployed today, free tier, queryable now",
             left=0.6, top=0.45, width=12.0, height=0.25,
             font_name=JBM_LIGHT, size=9, color=FADED_INK,
             spacing_hpt=120)

    # ── Left two-thirds: the manifest ──
    manifest_x = 0.6
    manifest_top = 1.1
    rows = [
        ("frontend",  "next.js 16  ·  vercel  ·  hunarmand.sal.lol"),
        ("backend",   "fastapi  ·  render  ·  alembic on every deploy"),
        ("ai core",   "fastapi  ·  hugging face spaces  ·  ed25519 sanad service"),
        ("postgres",  "neon  ·  pgvector  ·  shared between backend + ai core"),
        ("asr",       "bhashini  >  ai4bharat  >  groq  >  hf inference  >  whisper"),
        ("llm",       "openrouter  ·  free + premium models behind one interface"),
        ("embeddings","sentence-transformers  ·  intfloat / multilingual-e5-small"),
        ("signature", "ed25519 over rfc 8785 jcs  ·  jws-compact qr  ·  offline verify"),
    ]
    row_h = 0.42
    for i, (key, val) in enumerate(rows):
        ry = manifest_top + i * row_h
        # key column
        add_text(slide, key,
                 left=manifest_x, top=ry, width=1.6, height=0.32,
                 font_name=JBM_REG, size=11, bold=True, color=WALNUT,
                 spacing_hpt=60)
        # equals separator
        add_text(slide, "·",
                 left=manifest_x + 1.65, top=ry, width=0.2, height=0.32,
                 font_name=JBM_LIGHT, size=11, color=WALNUT_40)
        # value column
        add_text(slide, val,
                 left=manifest_x + 1.95, top=ry, width=7.5, height=0.32,
                 font_name=JBM_LIGHT, size=11, color=FADED_INK)

    # Closing line under the manifest
    closing_y = manifest_top + len(rows) * row_h + 0.25
    add_text(slide,
             "every dependency above is on a free tier today.",
             left=manifest_x, top=closing_y, width=10.0, height=0.32,
             font_name=JBM_REG, size=12, bold=True, color=WALNUT)

    # ── Right column: the live numbers ──
    nums_x = 10.4
    nums_top = 1.1
    nums = [
        ("18",   "seeded\nartisan masters"),
        ("54",   "bookable\nworkshops"),
        ("2.8s", "median latency,\nsign + persist"),
    ]
    spacing = 1.55
    for i, (n, cap) in enumerate(nums):
        y = nums_top + i * spacing
        add_text(slide, n,
                 left=nums_x, top=y, width=2.5, height=0.95,
                 font_name=JBM_LIGHT, size=64, color=WALNUT,
                 anchor=MSO_ANCHOR.TOP)
        add_text(slide, cap,
                 left=nums_x, top=y + 0.95, width=2.5, height=0.55,
                 font_name=JBM_LIGHT, size=9, color=FADED_INK,
                 spacing_hpt=60, line_spacing=1.4)

    # Punch line at the bottom
    punch_y = SLIDE_H - 0.85
    add_disk(slide, cx=0.7, cy=punch_y + 0.11, diameter_pt=8, fill=ROUGE)
    add_text(slide,
             "every number is queryable from the public api right now.",
             left=0.95, top=punch_y, width=10.5, height=0.32,
             font_name=JBM_REG, size=12, bold=True, color=WALNUT)
    add_line(slide,
             x1=0.95, y1=punch_y + 0.34, x2=0.95 + 2.4, y2=punch_y + 0.34,
             color=ROUGE, weight_pt=0.9)

    add_footer(slide, 5)


# ─────────────────────────────────────────────────────────────────────
#  (legacy, retained for reference) the old per-slide builders
# ─────────────────────────────────────────────────────────────────────
def slide_insight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Caption above glyph (left)
    add_text(
        slide,
        "talim, kanihama, c. 1820",
        left=0.6, top=0.55, width=4.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )

    # Glyph stack: 8 cells, each cell is two short angled walnut strokes
    # Cell: 0.4" tall × 0.3" wide. Centered around x = 1.55"
    glyph_x = 1.55
    glyph_top = 1.0
    cell_h = 0.55
    cell_w = 0.34
    import random
    random.seed(7)  # reproducible glyph
    for i in range(8):
        cy = glyph_top + i * cell_h
        # Two short strokes per cell, slightly different angles
        # Stroke A: top-left to mid-right
        ax1 = glyph_x - cell_w / 2 + random.uniform(0.02, 0.06)
        ay1 = cy + random.uniform(0.04, 0.10)
        ax2 = glyph_x + random.uniform(-0.02, 0.06)
        ay2 = cy + cell_h * 0.45 + random.uniform(-0.04, 0.04)
        add_line(slide, x1=ax1, y1=ay1, x2=ax2, y2=ay2,
                 color=WALNUT, weight_pt=0.8)
        # Stroke B: small descender
        bx1 = glyph_x + random.uniform(-0.03, 0.05)
        by1 = cy + cell_h * 0.40 + random.uniform(-0.03, 0.04)
        bx2 = glyph_x + cell_w / 2 - random.uniform(0.02, 0.07)
        by2 = cy + cell_h * 0.85 + random.uniform(-0.03, 0.04)
        add_line(slide, x1=bx1, y1=by1, x2=bx2, y2=by2,
                 color=WALNUT, weight_pt=0.8)

    glyph_bottom_y = glyph_top + 8 * cell_h  # ~ 5.4"

    # Caption above JSON (right)
    json_x = 7.4
    add_text(
        slide,
        "craft dna, hunarmand, 2026",
        left=json_x, top=0.55, width=5.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )

    # JSON body — split each line into key + value runs so we can colour
    # them differently.
    json_top = 1.0
    json_lines = [
        ("{", None),
        ("  \"piece\":      ",      "\"kani-buti pashmina shawl\","),
        ("  \"warp\":       ",      "\"single, 28 ends per cm\","),
        ("  \"weft\":       ",      "[\"pashmina, ladakh\", \"silk, varanasi\"],"),
        ("  \"knot\":       ",      "\"1940s srinagar twill-tapestry\","),
        ("  \"dye\":        ",      "[\"walnut bark\", \"saffron stigma\"],"),
        ("  \"kid\":        ",      "\"mohammad-yusuf:1\","),
        ("  \"signature\":  ",      "\"ed25519, jcs canonical\""),
        ("}", None),
    ]
    box = slide.shapes.add_textbox(
        Inches(json_x), Inches(json_top),
        Inches(SLIDE_W - json_x - 0.6), Inches(4.6),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    disable_autosize(tf)

    for i, (left, right) in enumerate(json_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = 1.45
        r1 = para.add_run()
        r1.text = left
        r1.font.name = JBM_REG
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = WALNUT
        if right is not None:
            r2 = para.add_run()
            r2.text = right
            r2.font.name = JBM_REG
            r2.font.size = Pt(11)
            r2.font.color.rgb = FADED_INK

    json_bottom_y = json_top + 4.4

    # Connecting rouge S-curve from glyph bottom to JSON top — approx as
    # two short hairline segments meeting in the middle of the slide.
    midx = (glyph_x + json_x) / 2
    midy = (glyph_bottom_y + json_top) / 2
    add_line(
        slide,
        x1=glyph_x + 0.05, y1=glyph_bottom_y - 0.1,
        x2=midx,           y2=midy + 0.4,
        color=ROUGE, weight_pt=1.0,
    )
    add_line(
        slide,
        x1=midx,           y1=midy + 0.4,
        x2=json_x - 0.08,  y2=json_top + 0.55,
        color=ROUGE, weight_pt=1.0,
    )

    # Closing line, full-width
    add_text(
        slide,
        "same instruction.  seventh medium.",
        left=0.6, top=SLIDE_H - 1.4, width=SLIDE_W - 1.2, height=0.4,
        font_name=JBM_REG, size=18, bold=True, color=WALNUT,
    )

    add_footer(slide, 3)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 04 — the system (diagonal pipeline)
# ─────────────────────────────────────────────────────────────────────
def slide_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top-left caption
    add_text(
        slide,
        "architecture / 2026.q2",
        left=0.6, top=0.45, width=5.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    # Nodes
    nodes = [
        # (x, y, diameter_pt, fill, label, function)
        (1.2,  1.6,  10, WALNUT, "vault",
         "captures what the master never wrote down"),
        (4.4,  3.1,  12, ROUGE,  "sanad",
         "signs the piece against the master's keypair"),
        (7.6,  4.55, 10, WALNUT, "ustaad",
         "sells the master's own time, not the middleman's"),
        (10.4, 6.0,  10, WALNUT, "bazaar",
         "sells the object, with provenance attached"),
    ]

    # Connectors between nodes (with disk margins)
    for i in range(len(nodes) - 1):
        x1, y1, d1, *_ = nodes[i]
        x2, y2, d2, *_ = nodes[i + 1]
        d1_in = d1 / 72.0
        d2_in = d2 / 72.0
        # Vector from node1 to node2
        dx = x2 - x1
        dy = y2 - y1
        import math
        L = math.hypot(dx, dy)
        ux, uy = dx / L, dy / L
        sx, sy = x1 + ux * d1_in / 2, y1 + uy * d1_in / 2
        ex, ey = x2 - ux * d2_in / 2, y2 - uy * d2_in / 2
        add_line(
            slide,
            x1=sx, y1=sy, x2=ex, y2=ey,
            color=WALNUT, weight_pt=0.9,
        )

    # Disks + labels
    for x, y, d, fill, label, fn in nodes:
        add_disk(slide, cx=x, cy=y, diameter_pt=d, fill=fill)
        label_x = x + (d / 72.0) / 2 + 0.22
        # Caption widths must not run off the slide. Cap at the
        # remaining horizontal room minus a 0.6" right margin.
        avail_w = max(2.0, SLIDE_W - label_x - 0.6)
        add_text(
            slide,
            label,
            left=label_x, top=y - 0.32, width=avail_w, height=0.32,
            font_name=JBM_REG, size=16, bold=True, color=WALNUT,
        )
        add_text(
            slide,
            fn,
            left=label_x, top=y - 0.04, width=avail_w, height=0.32,
            font_name=JBM_LIGHT, size=10, color=FADED_INK,
            spacing_hpt=60,
        )

    # The keystone treatment sits in the upper-right white space, well
    # clear of the vault description and the diagonal connector. Anchored
    # right of the slide so it reads as a margin gloss, not as a label
    # competing with the pipeline.
    sanad_x, sanad_y = nodes[1][0], nodes[1][1]
    keystone_left = 9.6
    keystone_y = 0.95
    add_line(
        slide,
        x1=keystone_left, y1=keystone_y,
        x2=keystone_left + 1.4, y2=keystone_y,
        color=ROUGE, weight_pt=1.0,
    )
    add_text(
        slide,
        "the keystone",
        left=keystone_left, top=keystone_y + 0.08, width=2.6, height=0.32,
        font_name=JBM_REG, size=11, bold=True, color=WALNUT,
    )
    add_text(
        slide,
        "provenance is what buyers pay\na premium for, today",
        left=keystone_left, top=keystone_y + 0.36, width=3.3, height=0.7,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60, line_spacing=1.4,
    )

    # A faint walnut leader line drops from the keystone label down to
    # the sanad rouge node, so the gloss stays anchored to the keystone.
    add_line(
        slide,
        x1=keystone_left + 0.1, y1=keystone_y + 0.85,
        x2=sanad_x + 0.18,      y2=sanad_y - 0.12,
        color=WALNUT_30, weight_pt=0.5, dash="sysDot",
    )

    add_footer(slide, 4)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 05 — the economics
# ─────────────────────────────────────────────────────────────────────
def slide_economics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(
        slide,
        "who keeps the value of a hand-knotted shawl",
        left=0.6, top=0.55, width=8.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    bar_x = 2.0
    bar_w = 9.0
    bar_h = 0.45

    # ── Today bar ──
    top1 = 2.4
    # left label
    add_text(
        slide,
        "today",
        left=0.6, top=top1 + (bar_h - 0.22) / 2, width=1.3, height=0.32,
        font_name=JBM_REG, size=11, bold=True, color=WALNUT,
    )
    # walnut segment (70%)
    seg1_w = bar_w * 0.78
    add_filled_rect(slide, left=bar_x, top=top1, width=seg1_w, height=bar_h,
                    fill=WALNUT)
    # faded segment (22% — middle of 15-30 range)
    seg2_w = bar_w - seg1_w
    add_filled_rect(slide, left=bar_x + seg1_w, top=top1, width=seg2_w, height=bar_h,
                    fill=FADED_35)

    # captions under top bar
    cap_top = top1 + bar_h + 0.12
    add_text(
        slide,
        "middleman chain, 70 to 85 percent",
        left=bar_x, top=cap_top, width=seg1_w, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )
    add_text(
        slide,
        "artisan, 15 to 30 percent",
        left=bar_x + seg1_w, top=cap_top, width=seg2_w + 1.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )

    # ── Hunarmand bar ──
    top2 = top1 + bar_h + 0.6 + 0.5
    # left label
    add_text(
        slide,
        "hunarmand",
        left=0.6, top=top2 + (bar_h - 0.22) / 2, width=1.3, height=0.32,
        font_name=JBM_REG, size=11, bold=True, color=ROUGE,
    )
    # rouge segment (8%)
    seg1b_w = bar_w * 0.08
    add_filled_rect(slide, left=bar_x, top=top2, width=seg1b_w, height=bar_h,
                    fill=ROUGE)
    # walnut segment (92%)
    seg2b_w = bar_w - seg1b_w
    add_filled_rect(slide, left=bar_x + seg1b_w, top=top2, width=seg2b_w, height=bar_h,
                    fill=WALNUT)

    # captions under hunarmand bar. Both anchored at segment edges so
    # the eye reads them as labels of the slivers above. The rouge
    # segment is narrow enough that its caption right-aligns to end at
    # the segment edge, leaving comfortable space before the artisan
    # caption that left-aligns to start at the same edge.
    cap2_top = top2 + bar_h + 0.12
    add_text(
        slide,
        "platform, 6 to 9 percent",
        left=bar_x - 1.6, top=cap2_top, width=1.6 + seg1b_w - 0.05, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        align=PP_ALIGN.RIGHT, spacing_hpt=60,
    )
    add_text(
        slide,
        "artisan, 91 to 94 percent",
        left=bar_x + seg1b_w + 0.15, top=cap2_top, width=seg2b_w, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )

    # Closing line
    add_text(
        slide,
        "the artisan keeps the rest.",
        left=bar_x, top=cap2_top + 0.7, width=10.0, height=0.4,
        font_name=JBM_REG, size=14, bold=True, color=WALNUT,
    )

    add_footer(slide, 5)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 06 — the proof
# ─────────────────────────────────────────────────────────────────────
def slide_proof(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Top caption
    add_text(
        slide,
        "live build / hunarmand.sal.lol",
        left=0.6, top=0.55, width=6.0, height=0.25,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=120,
    )

    rows = [
        ("18",    "seeded artisan masters, four crafts, six villages"),
        ("54",    "bookable workshops, four formats, six regions"),
        ("2.8 s", "median latency, ed25519 sign and persist, neon to vercel"),
    ]

    # Rows: numbers 96pt Light walnut. Captions to the right, 11pt Light faded ink.
    row_top0 = 1.5
    row_step = 1.6  # vertical breathing
    for i, (number, caption) in enumerate(rows):
        y = row_top0 + i * row_step
        # Big number
        add_text(
            slide,
            number,
            left=0.6, top=y, width=4.0, height=1.4,
            font_name=JBM_LIGHT, size=96, color=WALNUT,
            anchor=MSO_ANCHOR.TOP,
        )
        # Caption right of number, vertically near cap height
        add_text(
            slide,
            caption,
            left=4.7, top=y + 0.55, width=8.0, height=0.4,
            font_name=JBM_LIGHT, size=11, color=FADED_INK,
            spacing_hpt=60,
        )

    # Bottom left line, sitting clear of the last "2.8 s" descender
    bl_y = SLIDE_H - 0.85
    add_text(
        slide,
        "nothing on these slides is a mockup.",
        left=0.6, top=bl_y, width=4.6, height=0.3,
        font_name=JBM_REG, size=11, bold=True, color=WALNUT,
    )
    # Rouge hairline under bold line
    add_line(
        slide,
        x1=0.6, y1=bl_y + 0.32,
        x2=0.6 + 2.0, y2=bl_y + 0.32,
        color=ROUGE, weight_pt=0.9,
    )
    # Rouge dot, then small caption to its right
    add_disk(slide, cx=5.5, cy=bl_y + 0.11, diameter_pt=8, fill=ROUGE)
    add_text(
        slide,
        "every number is queryable from the public api right now",
        left=5.7, top=bl_y, width=7.5, height=0.3,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        spacing_hpt=60,
    )

    add_footer(slide, 6)


# ─────────────────────────────────────────────────────────────────────
#  SLIDE 07 — close
# ─────────────────────────────────────────────────────────────────────
def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_parchment_bg(slide, SLIDE_W, SLIDE_H)

    # Quote: two lines, line 1 at 60% alpha, line 2 at 100% bold.
    # 28pt mono fits ~50 chars on a single line in a 12.3"-wide box and
    # has the weight a closing slide deserves on a 13.3" canvas.
    quote_left = 0.5
    quote_w = SLIDE_W - 1.0
    quote_top = 2.45

    # Line 1: faded (precomputed walnut at 60%)
    add_text(
        slide,
        "we are not bringing foreign technology to kashmir.",
        left=quote_left, top=quote_top, width=quote_w, height=0.7,
        font_name=JBM_LIGHT, size=28, color=WALNUT_60,
        align=PP_ALIGN.CENTER,
    )

    # Line 2: bold full walnut
    add_text(
        slide,
        "we are upgrading what kashmir already invented.",
        left=quote_left, top=quote_top + 0.78, width=quote_w, height=0.7,
        font_name=JBM_REG, size=28, bold=True, color=WALNUT,
        align=PP_ALIGN.CENTER,
    )

    # Rouge hairline below quote, centered
    rule_y = quote_top + 0.78 + 0.7 + 0.5
    rule_w = 1.4
    add_line(
        slide,
        x1=(SLIDE_W - rule_w) / 2, y1=rule_y,
        x2=(SLIDE_W + rule_w) / 2, y2=rule_y,
        color=ROUGE, weight_pt=0.9,
    )

    # Wordmark, centered
    add_text(
        slide,
        "hunarmand",
        left=0, top=rule_y + 0.3, width=SLIDE_W, height=0.4,
        font_name=JBM_REG, size=18, bold=True, color=WALNUT,
        align=PP_ALIGN.CENTER,
    )

    # Caption beneath wordmark
    add_text(
        slide,
        "tacit knowledge os, srinagar, 2026",
        left=0, top=rule_y + 0.78, width=SLIDE_W, height=0.3,
        font_name=JBM_LIGHT, size=9, color=FADED_INK,
        align=PP_ALIGN.CENTER, spacing_hpt=120,
    )

    add_footer(slide, 6)


# ─────────────────────────────────────────────────────────────────────
def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs)         # 01 — cover
    slide_disappearance(prs) # 02 — problem
    slide_solution(prs)      # 03 — solution
    slide_methodology(prs)   # 04 — methodology
    slide_technology(prs)    # 05 — technology
    slide_close(prs)         # 06 — close

    out = Path(__file__).with_name("hunarmand.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build()
    print(f"wrote {out}")
